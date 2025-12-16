#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Document generation module for Telegram bot
This module generates PowerPoint presentations and Word documents
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import random

# Sample content for documents
SAMPLE_SECTIONS = [
    "Kirish",
    "Asosiy qism",
    "Xulosa",
    "Foydalanilgan adabiyotlar",
    "Tahlil",
    "Natijalar",
    "Takliflar",
    "Muhokama",
    "Amaliy qism",
    "Nazariy asoslar"
]

SAMPLE_CONTENT = {
    "Kirish": [
        "Kirish qismida muammoning dolzarbligi, tadqiqot obyekti va predmeti, maqsad va vazifalar, tadqiqot usullari, ilmiy yangilik va amaliy ahamiyati bayon etiladi.",
        "Bu yerda tadqiqotning asosiy g'oyasi va uning dolzarbligi ta'kidlanadi.",
        "Kirish qismi tadqiqotning boshlanish nuqtasi hisoblanadi."
    ],
    "Asosiy qism": [
        "Asosiy qismda tadqiqot obyekti va predmeti, tadqiqot usullari, tahlil natijalari, amaliy qadamlar keltiriladi.",
        "Bu yerda muammoga chuqurroq kirib boriladi va uning asosiy jihatlari o'rganiladi.",
        "Asosiy qismda nazariy va amaliy ma'lumotlar tahlil qilinadi."
    ],
    "Xulosa": [
        "Xulosa qismida tadqiqot natijalari umumlashtiriladi, asosiy выводlar keltiriladi va kelajakdagi tadqiqotlar uchun yo'nalishlar aniqlanadi.",
        "Bu yerda tadqiqotning asosiy natijalari qisqacha bayon etiladi.",
        "Xulosa qismi tadqiqotning yakuni hisoblanadi."
    ],
    "Foydalanilgan adabiyotlar": [
        "1. O'zbekiston Respublikasi Prezidenti farmonlari va qarorlari",
        "2. O'zbekiston Respublikasi Vazirlar Mahkamasi qarorlari",
        "3. Ilmiy maqolalar va monografiyalar",
        "4. Internet manbalar"
    ],
    "Tahlil": [
        "Tahlil qismida yig'ilgan ma'lumotlar tahlil qilinadi, ularda qonuniyatlar va bog'liqliklar aniqlanadi.",
        "Bu yerda statistik ma'lumotlar, grafiklar va jadvallar keltiriladi.",
        "Tahlil qismi tadqiqotning eng muhim qismlaridan biri hisoblanadi."
    ],
    "Natijalar": [
        "Natijalar qismida tadqiqot natijalari keltiriladi, ularga tahlil qilinadi va xulosalar chiqariladi.",
        "Bu yerda tadqiqotning amaliy ahamiyati ta'kidlanadi.",
        "Natijalar qismi tadqiqotning yakuniy natijalarini o'z ichiga oladi."
    ],
    "Takliflar": [
        "Takliflar qismida tadqiqot natijalariga asoslanib amaliy takliflar keltiriladi.",
        "Bu yerda muammoni hal qilish yo'llari taklif qilinadi.",
        "Takliflar qismi tadqiqotning amaliy qadami hisoblanadi."
    ],
    "Muhokama": [
        "Muhokama qismida tadqiqot natijalari boshqa tadqiqotlar bilan solishtiriladi, ularda o'xshashliklar va farqlar aniqlanadi.",
        "Bu yerda tadqiqotning ilmiy yangiligi ta'kidlanadi.",
        "Muhokama qismi tadqiqotning ilmiy ahamiyatini oshiradi."
    ],
    "Amaliy qism": [
        "Amaliy qismda nazariy ma'lumotlar amaliyotda qanday qo'llanilishi ko'rsatiladi.",
        "Bu yerda amaliy misollar va vaziya tahlili keltiriladi.",
        "Amaliy qism tadqiqotning amaliy qadami hisoblanadi."
    ],
    "Nazariy asoslar": [
        "Nazariy asoslar qismida tadqiqotga asos bo'lgan nazariy ma'lumotlar keltiriladi.",
        "Bu yerda tadqiqotning nazariy asoslari bayon etiladi.",
        "Nazariy asoslar qismi tadqiqotning nazariy qadami hisoblanadi."
    ]
}

def generate_presentation(title, pages, output_path=None):
    """
    Generate PowerPoint presentation
    
    Args:
        title: Presentation title
        pages: Number of pages (as string: '5-10', '15-20', etc.)
        output_path: Output file path
    
    Returns:
        Path to generated presentation
    """
    # Create presentation
    prs = Presentation()
    
    # Set slide width and height
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = "Slaydtop Bot"
    
    # Determine number of content slides based on pages
    page_ranges = {
        '5-10': 3,
        '15-20': 5,
        '25-30': 7,
        '35-40': 9,
        '45-50': 11
    }
    
    num_slides = page_ranges.get(pages, 3)
    
    # Add content slides
    content_slide_layout = prs.slide_layouts[1]
    
    for i in range(num_slides):
        slide = prs.slides.add_slide(content_slide_layout)
        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1]
        
        # Choose random section
        section = random.choice(SAMPLE_SECTIONS)
        title_placeholder.text = section
        
        # Add content
        text_frame = body_placeholder.text_frame
        text_frame.text = random.choice(SAMPLE_CONTENT[section])
        
        # Add bullet points
        for j in range(2):
            p = text_frame.add_paragraph()
            p.text = random.choice(SAMPLE_CONTENT[section])
            p.level = 1
    
    # Save presentation
    if output_path is None:
        output_path = f"presentation_{title.replace(' ', '_')}.pptx"
    
    prs.save(output_path)
    
    return output_path

def generate_word_document(title, document_type, pages, output_path=None):
    """
    Generate Word document
    
    Args:
        title: Document title
        document_type: Type of document ('independent_work', 'reference', 'thesis', 'article')
        pages: Number of pages (as string: '5-10', '15-20', etc.)
        output_path: Output file path
    
    Returns:
        Path to generated document
    """
    # Create document
    doc = Document()
    
    # Set document margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.25)
        section.right_margin = Inches(1.25)
    
    # Add title
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Add author info
    author_para = doc.add_paragraph()
    author_para.add_run('Muallif: ').bold = True
    author_para.add_run('Slaydtop Bot')
    author_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Determine number of sections based on pages
    page_ranges = {
        '5-10': 3,
        '15-20': 5,
        '25-30': 7,
        '35-40': 9,
        '45-50': 11
    }
    
    num_sections = page_ranges.get(pages, 3)
    
    # Add sections
    for i in range(num_sections):
        section = random.choice(SAMPLE_SECTIONS)
        
        # Add section heading
        doc.add_heading(section, level=1)
        
        # Add content
        for j in range(3):
            para = doc.add_paragraph(random.choice(SAMPLE_CONTENT[section]))
            para.style = 'Body Text'
        
        # Add spacing
        doc.add_paragraph()
    
    # Add conclusion
    doc.add_heading('Xulosa', level=1)
    conclusion_para = doc.add_paragraph()
    conclusion_para.add_run('Xulosa: ')
    conclusion_para.add_run(random.choice(SAMPLE_CONTENT['Xulosa']))
    
    # Add references
    doc.add_heading('Foydalanilgan adabiyotlar', level=1)
    for ref in SAMPLE_CONTENT['Foydalanilgan adabiyotlar']:
        doc.add_paragraph(ref, style='List Number')
    
    # Save document
    if output_path is None:
        output_path = f"document_{title.replace(' ', '_')}.docx"
    
    doc.save(output_path)
    
    return output_path

def generate_presentation_from_file(title, file_url, pages, output_path=None):
    """
    Generate presentation from file or URL
    
    Args:
        title: Presentation title
        file_url: File path or URL
        pages: Number of pages
        output_path: Output file path
    
    Returns:
        Path to generated presentation
    """
    # For now, just generate a regular presentation
    # In real implementation, this would process the file/URL content
    return generate_presentation(title, pages, output_path)

# Example usage
if __name__ == '__main__':
    # Generate PowerPoint presentation
    ppt_path = generate_presentation(
        title="Test Presentation",
        pages='5-10',
        output_path='test_presentation.pptx'
    )
    print(f"Presentation generated: {ppt_path}")
    
    # Generate Word document
    doc_path = generate_word_document(
        title="Test Document",
        document_type='independent_work',
        pages='5-10',
        output_path='test_document.docx'
    )
    print(f"Document generated: {doc_path}")
